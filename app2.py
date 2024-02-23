from flask import Flask, render_template, request, send_from_directory, send_file
import os
from docx import Document
from pptx import Presentation
from PIL import Image
from io import BytesIO
import base64
import fitz
import yadisk

app = Flask(__name__)

# Токен для доступа к Яндекс.Диску
y = yadisk.YaDisk(token="y0_AgAAAAAl5jRpAAtKfAAAAAD6_t3BAABRqHACUIJMrJ4PRAckXHu7iVOojw")

UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def list_files_on_yadisk(folder_path):
    files = y.listdir(folder_path)
    files_list = [(file.name, file.path) for file in files]
    return files_list

def search_files(files, query):
    if not query:
        return files
    query = query.lower()
    return [file for file in files if query in file[0].lower()]

@app.route('/', methods=['GET', 'POST'])
def index():
    query = request.form.get('query')
    sort_by = request.form.get('sort_by')
    yadisk_files = list_files_on_yadisk('/uploads')

    # Применяем поиск
    yadisk_files = search_files(yadisk_files, query)

    # Сортируем файлы, если указан критерий сортировки
    if sort_by == 'name':
        yadisk_files.sort(key=lambda x: x[0])
    elif sort_by == 'date':
        yadisk_files.sort(key=lambda x: x[1].split('/')[-1])
        if filename.endswith('.docx'):
                text = process_word(file_path)
                return render_template('word_template.html', text=text, filename=filename) 
            
        elif filename.endswith('.pdf'):
                filename, images = process_pdf(file_path)
                return render_template('pdf_template.html', filename=filename, images=images)
            
        elif filename.endswith('.pptx'):
                slides_data = process_pptx(file_path)
                return render_template('pptx_template.html', slides=slides_data, filename=filename)
            
        else:
                text = "Unsupported file format"
                slides_data = [{"text": text, "images": []}]
                return render_template('pptx_template.html', slides=slides_data, filename=filename)
    return render_template('student.html', yadisk_files=yadisk_files)


@app.route('/upload', methods=['POST'])
def upload_file():
    if request.method == 'POST':
        file = request.files['file']
        if file:
            filename = file.filename
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            # Не загружаем файл на Яндекс.Диск, только сохраняем локально
            # y.upload(file_path, f'/uploads/{os.path.basename(file_path)}')

            if filename.endswith('.docx'):
                text = process_word(file_path)
                return render_template('word_template.html', text=text, filename=filename) 
            
            elif filename.endswith('.pdf'):
                filename, images = process_pdf(file_path)
                return render_template('pdf_template.html', filename=filename, images=images)
            
            elif filename.endswith('.pptx'):
                slides_data = process_pptx(file_path)
                return render_template('pptx_template.html', slides=slides_data, filename=filename)
            
            else:
                text = "Unsupported file format"
                slides_data = [{"text": text, "images": []}]
                return render_template('pptx_template.html', slides=slides_data, filename=filename)

    return "File upload failed."


@app.route('/open_from_yadisk/<path:file_path>')
def open_from_yadisk(file_path):
    temp_file = f"temp_{os.path.basename(file_path)}"
    with open(temp_file, "wb") as f:
        y.download(file_path, f)

    return send_file(temp_file, as_attachment=True)
def process_word(file_path):
    document = Document(file_path)
    text = ""
    for paragraph in document.paragraphs:
        text += paragraph.text + "\n"
    return text

def process_pdf(file_path):
    images = []
    filename = os.path.basename(file_path)
    with fitz.open(file_path) as pdf:
        for page in pdf:
            pixmap = page.get_pixmap(alpha=True)
            img = Image.frombytes("RGBA", [pixmap.width, pixmap.height], pixmap.samples)
            image_stream = BytesIO()
            img.save(image_stream, format="PNG")
            image_base64 = base64.b64encode(image_stream.getvalue()).decode('utf-8')
            images.append(image_base64)
    return filename, images

def process_pptx(file_path):
    presentation = Presentation(file_path)
    slides_data = []
    for slide in presentation.slides:
        slide_data = {"text": "", "images": []}
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_data["text"] += shape.text + "\n"
            elif shape.shape_type == 13:
                image_stream = BytesIO()
                img = shape.image
                img_bytes = img.blob
                image_stream.write(img_bytes)
                image_data = base64.b64encode(image_stream.getvalue()).decode('utf-8')
                slide_data["images"].append(image_data)
        slides_data.append(slide_data)
    return slides_data
if __name__ == '__main__':
    app.run(debug=True, port=5001)

